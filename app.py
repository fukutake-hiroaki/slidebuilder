import streamlit as st
import openai
import json
import os
from pptx import Presentation
from pptx.util import Pt
import io

# ディレクトリ設定
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MANUAL_PATH = os.path.join(BASE_DIR, 'slide manual.jsonl')
MASTERS_DIR = os.path.join(BASE_DIR, 'masters')
OUTPUTS_DIR = os.path.join(BASE_DIR, 'outputs')
UPLOADS_DIR = os.path.join(BASE_DIR, 'uploads')

# マニュアル読み込み
def load_manual():
    manuals = []
    with open(MANUAL_PATH, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                manuals.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    return manuals

st.title('スライド自動生成アプリ')

# OpenAI APIキーの取得（API KEY.txtから読み込み）
API_KEY_FILE = os.path.join(BASE_DIR, 'API KEY.txt')
openai_api_key = None
if os.path.exists(API_KEY_FILE):
    with open(API_KEY_FILE, 'r', encoding='utf-8') as f:
        openai_api_key = f.read().strip()
else:
    openai_api_key = st.text_input('OpenAI APIキーを入力してください', type='password')
if openai_api_key:
    import openai as openai_lib
    client = openai_lib.OpenAI(api_key=openai_api_key)

# テキスト入力欄（長文対応）
text = st.text_area('スライド化したいテキストを入力してください', height=400)
if text:
    with open(MANUAL_PATH, 'r', encoding='utf-8') as f:
        manual_text = f.read()
    master_path = os.path.join(BASE_DIR, 'slide.pptx')
    prs = Presentation(master_path)
    layout_names = [l.name for l in prs.slide_layouts]

    if openai_api_key:
        if 'slide_plan' not in st.session_state:
            st.session_state['slide_plan'] = ''
        if st.button('スライド構成案を作成する'):
            with st.spinner('スライド構成案を生成中...'):
                prompt = f"""
あなたは優秀なスライド変換AIです。以下のテキストを、スライド形式に変換します。スライドタイトルと各スライドの内容を、情報量を極力減らさずに、日本語で出力してください。スライドは10枚程度で、スライドの内容は元の情報を適切に残し、スライド1枚あたりの情報量を多くしてください。要約は極力しないでください。

【追加指示】
・各スライドごとに、どのスライドレイアウト（スライドマニュアルのlayout_name）を使うべきかも必ず明記してください。
・マニュアルを確認し、中見出し、本文、2カラム、表など、どのレイアウトが適切かを考えて割り当ててください。
・表が必要な場合は「表を挿入する」旨も明記してください。
・備考や注意点を出力しないでください。

[スライドマニュアル]
{manual_text}

---
{text}
---
"""
                try:
                    completion = client.chat.completions.create(
                        model="gpt-4.1",
                        messages=[{"role": "user", "content": prompt}]
                    )
                    st.session_state['slide_plan'] = completion.choices[0].message.content
                except Exception as e:
                    st.error(f"OpenAI APIエラー: {e}")
                    st.session_state['slide_plan'] = ''
        if st.session_state['slide_plan']:
            st.subheader('スライド構成案')
            edited_plan = st.text_area('スライド構成案を編集できます', value=st.session_state['slide_plan'], height=300)
            if st.button('この構成案を承認してスライド作成'):
                with st.spinner('スライド詳細案を作成中...'):
                    detail_prompt = f"""
あなたは優秀なスライド作成AIです。以下の条件を厳守してください。
・レイアウトは、スライドマニュアルのlayout_nameからユースケースに応じて選ぶこと。
・文章は、スライドマニュアルを参照し、該当するレイアウトのplaceholderに対して過不足なく入力すること。
・boxesのキーは必ずplaceholderのidx（数値。例: 0, 10, 11 など）で指定してください。BODY-10やTITLE-0のような文字列ではなく、"10"や"0"のような数字の文字列で指定してください。
・もしスライド内に表（テーブル）が必要な場合は、"tables"キーを追加し、各表ごとに"idx"（挿入先placeholderのidx、数値）と"data"（2次元配列の表データ）を含めてください。例: "tables": [{{"idx": 11, "data": [["見出し1", "見出し2"], ["値1", "値2"]]}}]

以下のスライド構成案とスライドマニュアルを参考に、各スライドごとに「レイアウト名」「テキストボックスidx」「テキストボックスの内容」「必要に応じて表データ」をJSON形式で出力してください。JSONの各要素は1スライドに対応し、"layout"（レイアウト名）、"boxes"（テキストボックスidxと内容の辞書）、"tables"（必要な場合のみ）を含めてください。余計な説明や文章は不要です。

[スライド構成案]
{edited_plan}

[スライドマニュアル]
{manual_text}
"""
                    try:
                        detail_completion = client.chat.completions.create(
                            model="gpt-4.1",
                            messages=[{"role": "user", "content": detail_prompt}]
                        )
                        detailed_json = detail_completion.choices[0].message.content
                    except Exception as e:
                        st.error(f"OpenAI APIエラーまたは詳細案生成エラー: {e}")
                        detailed_json = None
                if detailed_json:
                    try:
                        slides_data = json.loads(detailed_json)
                    except Exception:
                        try:
                            slides_data = json.loads(detailed_json[detailed_json.find('['):detailed_json.rfind(']')+1])
                        except Exception as e_json2:
                            st.error(f"詳細案JSONパース失敗: {e_json2}")
                            import traceback
                            st.error(traceback.format_exc())
                            slides_data = None
                    if slides_data:
                        master_path = os.path.join(BASE_DIR, 'slide.pptx')
                        prs = Presentation(master_path)
                        layout_name_to_idx = {l.name: i for i, l in enumerate(prs.slide_layouts)}
                        layout_names = list(layout_name_to_idx.keys())
                        for idx, slide_info in enumerate(slides_data):
                            layout = slide_info.get('layout')
                            if layout not in layout_name_to_idx:
                                st.warning(f"スライド{idx+1}: layout名 '{layout}' がpptxテンプレートに存在しません。既存レイアウト名: {layout_names}。デフォルト(1番)を使用します。")
                            layout_idx = layout_name_to_idx.get(layout, 1)
                            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                            box_idxs = list(slide_info.get('boxes', {}).keys())
                            placeholder_idxs = [ph.placeholder_format.idx for ph in slide.placeholders]
                            # テキスト挿入
                            for box_idx, box_content in slide_info.get('boxes', {}).items():
                                try:
                                    target_idx = int(box_idx)
                                except Exception:
                                    st.warning(f"スライド{idx+1}: boxのキー '{box_idx}' はintに変換できません。boxesのキーは数値(idx)である必要があります。")
                                    continue
                                found = False
                                for shape in slide.placeholders:
                                    if shape.placeholder_format.idx == target_idx:
                                        shape.text = str(box_content)
                                        found = True
                                        break
                                if not found:
                                    st.warning(f"スライド{idx+1}: プレースホルダーidx '{target_idx}' が見つかりません。利用可能なidx: {placeholder_idxs}")
                            # 表挿入
                            tables = slide_info.get('tables', [])
                            for table_info in tables:
                                table_idx = table_info.get('idx')
                                data = table_info.get('data')
                                if table_idx is None or data is None:
                                    st.warning(f"スライド{idx+1}: tablesの要素に'idx'または'data'がありません: {table_info}")
                                    continue
                                try:
                                    table_idx = int(table_idx)
                                except Exception:
                                    st.warning(f"スライド{idx+1}: tablesの'idx'がintに変換できません: {table_info}")
                                    continue
                                found = False
                                for shape in slide.placeholders:
                                    if shape.placeholder_format.idx == table_idx:
                                        left = shape.left
                                        top = shape.top
                                        width = shape.width
                                        height = shape.height
                                        rows = len(data)
                                        cols = len(data[0]) if rows > 0 else 0
                                        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                                        table = table_shape.table
                                        for r in range(rows):
                                            for c in range(cols):
                                                table.cell(r, c).text = str(data[r][c])
                                        found = True
                                        break
                                if not found:
                                    st.warning(f"スライド{idx+1}: 表用のプレースホルダーidx '{table_idx}' が見つかりません。利用可能なidx: {placeholder_idxs}")
                        pptx_bytes = io.BytesIO()
                        prs.save(pptx_bytes)
                        pptx_bytes.seek(0)
                        st.success('pptxスライドを生成しました')
                        st.download_button('pptxをダウンロード', pptx_bytes, file_name='generated_slide.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    else:
        st.info('OpenAI APIキーを入力してください')
